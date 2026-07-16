"""Extract text, tables, notes, and images from every PPTX in Lecture presentations and the labs folder."""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from config import EXTRACTED, LECTURES_DIR, LABS_DIR, slugify


def extract_shape_text(shape, slide_texts):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs) or para.text
            if text and text.strip():
                slide_texts.append(text.strip())
    if shape.has_table:
        table = shape.table
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        slide_texts.append({"table": rows})
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            extract_shape_text(sub, slide_texts)


def extract_one(pptx_path: Path, out_dir: Path, category: str) -> dict:
    file_id = slugify(pptx_path.name)
    img_dir = EXTRACTED / "images" / file_id
    img_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    slides_out = []
    total_images = 0
    for idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        images = []
        notes = ""
        for shape in slide.shapes:
            extract_shape_text(shape, slide_texts)
            if shape.shape_type == 13:  # PICTURE
                try:
                    image = shape.image
                    ext = image.ext
                    img_path = img_dir / f"slide_{idx:03d}_img{total_images + 1}.{ext}"
                    img_path.write_bytes(image.blob)
                    images.append(str(img_path.relative_to(EXTRACTED.parent)))
                    total_images += 1
                except Exception as e:
                    images.append(f"ERROR extracting image: {e}")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides_out.append({
            "slideNumber": idx,
            "texts": slide_texts,
            "notes": notes,
            "images": images,
        })

    result = {
        "id": file_id,
        "sourcePath": str(pptx_path),
        "category": category,
        "fileName": pptx_path.name,
        "slideCount": len(prs.slides.__iter__().__length_hint__()) if hasattr(prs.slides, "__length_hint__") else len(list(prs.slides)),
        "slides": slides_out,
        "totalImages": total_images,
    }
    result["slideCount"] = len(slides_out)

    out_path = out_dir / f"{file_id}.json"
    out_path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    return {
        "id": file_id,
        "sourcePath": str(pptx_path),
        "outputPath": str(out_path),
        "slideCount": len(slides_out),
        "totalImages": total_images,
        "status": "extracted",
    }


def main():
    manifest_entries = []
    jobs = []
    for f in sorted(LECTURES_DIR.glob("*.pptx")):
        jobs.append((f, EXTRACTED / "lectures", "lecture"))
    for f in sorted(LABS_DIR.glob("*.pptx")):
        jobs.append((f, EXTRACTED / "labs", "lab"))

    for pptx_path, out_dir, category in jobs:
        try:
            entry = extract_one(pptx_path, out_dir, category)
            print(f"OK  {category:8s} {pptx_path.name} -> {entry['slideCount']} slides, {entry['totalImages']} images")
        except Exception as e:
            entry = {
                "id": slugify(pptx_path.name),
                "sourcePath": str(pptx_path),
                "status": "error",
                "error": str(e),
            }
            print(f"ERR {category:8s} {pptx_path.name}: {e}", file=sys.stderr)
        manifest_entries.append(entry)

    manifest_path = EXTRACTED / "pptx_manifest.json"
    manifest_path.write_text(json.dumps(manifest_entries, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"\nWrote {manifest_path} with {len(manifest_entries)} entries")


if __name__ == "__main__":
    main()
